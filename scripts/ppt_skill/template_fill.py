"""
Template-based PPT filler - fills content into existing template,
preserving master slides, layouts, backgrounds, and visual design.

Usage:
    from template_fill import TemplatePptBuilder
    builder = TemplatePptBuilder("template.pptx", slides_data, images_dict)
    builder.build()
    builder.save("output.pptx")
"""

import os
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn


# ─── Slide layout matching heuristics ────────────────────
# Maps logical page types to layout name keywords (lowercase).
# The first matching layout index is used.

LAYOUT_MAP = {
    "cover":       ["title slide", "title", "cover", "\u5c01\u9762", "\u6807\u9898\u5e7b\u706f\u7247"],
    "section":     ["section header", "section", "\u8282\u6807\u9898", "\u7ae0\u8282"],
    "content":     ["title and content", "object", "content", "\u6807\u9898\u548c\u5185\u5bb9", "\u5185\u5bb9"],
    "two_content": ["two content", "comparison", "two objects", "\u4e24\u680f\u5185\u5bb9", "\u5bf9\u6bd4"],
    "blank":       ["blank", "\u7a7a\u767d"],
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _value_font_size(text: str, max_size: int = 44) -> int:
    """Choose a readable big-number font size that avoids wrapping into labels."""
    length = len(str(text or "").replace(" ", ""))
    if length <= 4:
        return max_size
    if length <= 7:
        return min(max_size, 38)
    if length <= 10:
        return min(max_size, 30)
    if length <= 14:
        return min(max_size, 25)
    return 20


def _set_ea_font(run, ea_font: str = "\u5fae\u8f6f\u96c5\u9ed1"):
    """Set East-Asian font for a run so Chinese text renders correctly."""
    if not ea_font:
        return
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', ea_font)


def _find_layout_index(prs, page_type: str) -> int:
    """Find the best-matching slide layout index for a given page type."""
    candidates = LAYOUT_MAP.get(page_type, LAYOUT_MAP["content"])
    for idx, layout in enumerate(prs.slide_layouts):
        name = (layout.name or "").lower()
        for keyword in candidates:
            if keyword in name:
                return idx
    # Fallback heuristics
    if page_type == "cover":
        return 0
    if page_type in ("next_steps",):
        return max(0, len(prs.slide_layouts) - 1)
    return min(1, len(prs.slide_layouts) - 1)


class TemplatePptBuilder:
    """Build PPT by filling content into an existing template.

    Preserves all master slides, layouts, backgrounds, and theme —
    only adds new slides and replaces placeholder text.
    """

    def __init__(self, template_path: str, slides_data: list, images: dict = None):
        self.template_path = template_path
        self.slides_data = slides_data
        self.images = images or {}
        self.prs = Presentation(template_path)
        self._clear_existing_slides()

    # ─── helpers ────────────────────────────────────────

    def _clear_existing_slides(self):
        """Delete all existing slides while keeping master & layouts intact."""
        while len(self.prs.slides) > 0:
            sld_id_lst = self.prs.slides._sldIdLst
            ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            rId = sld_id_lst[0].get(f'{{{ns}}}id') or sld_id_lst[0].get('r:id')
            if rId:
                try:
                    self.prs.part.drop_rel(rId)
                except Exception:
                    pass
            del sld_id_lst[0]

    def _add_slide(self, layout_idx: int):
        return self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

    def _fill_ph(self, slide, ph_type, text: str):
        """Fill the first placeholder of *ph_type* with *text*. Returns True if found."""
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = text
                if p.runs:
                    _set_ea_font(p.runs[0])
                return True
        return False

    def _fill_ph_multiline(self, slide, ph_type, lines: list):
        """Fill placeholder with multiple lines (one paragraph each)."""
        for shape in slide.placeholders:
            if shape.placeholder_format.type == ph_type:
                tf = shape.text_frame
                tf.clear()
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    if p.runs:
                        _set_ea_font(p.runs[0])
                return True
        return False

    def _add_tb(self, slide, left, top, width, height, text,
                font_size=14, bold=False, color="#333333",
                alignment=PP_ALIGN.LEFT):
        """Add a free-floating textbox to the slide."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = _hex_to_rgb(color)
        p.alignment = alignment
        if p.runs:
            _set_ea_font(p.runs[0])
        return txBox

    def _add_multiline_tb(self, slide, left, top, width, height, lines,
                          font_size=12, color="#333333"):
        """Add a textbox with multiple paragraphs."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = _hex_to_rgb(color)
            p.space_after = Pt(4)
        return txBox

    # ─── per-page builders ──────────────────────────────

    def build(self):
        for slide_data in self.slides_data:
            page_type = slide_data.get("type", "cover")
            method_name = f"_build_{page_type}"
            method = getattr(self, method_name, self._build_content)
            method(slide_data)

    def _build_cover(self, data):
        idx = _find_layout_index(self.prs, "cover")
        slide = self._add_slide(idx)

        title = data.get("title", "")
        subtitle = data.get("subtitle", "")

        t_ok = self._fill_ph(slide, PP_PLACEHOLDER.TITLE, title)
        if not t_ok:
            t_ok = self._fill_ph(slide, PP_PLACEHOLDER.CENTER_TITLE, title)
        s_ok = self._fill_ph(slide, PP_PLACEHOLDER.SUBTITLE, subtitle)

        if not t_ok and title:
            self._add_tb(slide, 1.2, 2.5, 10.9, 1.2, title, 36, True)
        if not s_ok and subtitle:
            self._add_tb(slide, 1.2, 3.6, 10.9, 0.6, subtitle, 18)

    def _build_problem(self, data):
        self._build_content(data)

    def _build_solution(self, data):
        """Solution overview — prefer TWO_CONTENT layout for left/right columns."""
        idx = _find_layout_index(self.prs, "two_content")
        slide = self._add_slide(idx)

        title = data.get("title", "")
        items = data.get("items", [])

        self._fill_ph(slide, PP_PLACEHOLDER.TITLE, title) or \
            self._fill_ph(slide, PP_PLACEHOLDER.CENTER_TITLE, title)

        problems = []
        solutions = []
        for item in items:
            if isinstance(item, dict):
                problems.append(item.get("problem", ""))
                solutions.append(item.get("solution", ""))

        # Try two body placeholders (left / right)
        body_idx = 0
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                tf = shape.text_frame
                tf.clear()
                if body_idx == 0:
                    tf.paragraphs[0].text = "\u5f53\u524d\u759b\u70b9"
                    for pt in problems:
                        p = tf.add_paragraph()
                        p.text = u"\u2022 " + pt
                else:
                    tf.paragraphs[0].text = "\u89e3\u51b3\u65b9\u6848"
                    for st in solutions:
                        p = tf.add_paragraph()
                        p.text = u"\u2022 " + st
                body_idx += 1

        # Fallback if no two-column placeholders
        if body_idx == 0 and problems:
            left = "\u3010\u759b\u70b9\u3011\n" + "\n".join(u"\u2022 " + p for p in problems)
            right = "\u3010\u65b9\u6848\u3011\n" + "\n".join(u"\u2022 " + s for s in solutions)
            self._add_multiline_tb(slide, 0.7, 1.3, 5.5, 5.0, left.split("\n"), 12)
            self._add_multiline_tb(slide, 6.5, 1.3, 5.5, 5.0, right.split("\n"), 12)

        tagline = data.get("tagline", "")
        if tagline:
            self._add_tb(slide, 0.7, 6.3, 11.9, 0.5, tagline, 14, True, "#4472C4")

    def _build_capabilities(self, data):
        idx = _find_layout_index(self.prs, "content")
        slide = self._add_slide(idx)

        title = data.get("title", "")
        self._fill_ph(slide, PP_PLACEHOLDER.TITLE, title) or \
            self._fill_ph(slide, PP_PLACEHOLDER.CENTER_TITLE, title)

        items = data.get("items", [])
        if not items:
            return

        # Use body placeholder if available
        body_used = False
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                tf = shape.text_frame
                tf.clear()
                for i, item in enumerate(items):
                    if isinstance(item, dict):
                        name = item.get("name", "")
                        desc = item.get("desc", "")
                    else:
                        name = str(item)
                        desc = ""
                    if i > 0:
                        tf.add_paragraph()
                    p = tf.add_paragraph()
                    run_b = p.add_run()
                    run_b.text = name
                    run_b.font.bold = True
                    run_b.font.size = Pt(14)
                    if desc:
                        p2 = tf.add_paragraph()
                        p2.text = desc
                        p2.font.size = Pt(12)
                body_used = True
                break

        # Else lay out as cards manually
        if not body_used:
            cols = 3
            for i, item in enumerate(items):
                if isinstance(item, dict):
                    name = item.get("name", "")
                    desc = item.get("desc", "")
                else:
                    name = str(item)
                    desc = ""
                col = i % cols
                row = i // cols
                x = 0.7 + col * 4.0
                y = 1.4 + row * 2.5
                self._add_tb(slide, x, y, 3.5, 0.5, name, 15, True, "#4472C4")
                self._add_tb(slide, x, y + 0.55, 3.5, 1.6, desc, 12, False, "#555555")

    def _build_how_it_works(self, data):
        idx = _find_layout_index(self.prs, "content")
        slide = self._add_slide(idx)

        title = data.get("title", "")
        self._fill_ph(slide, PP_PLACEHOLDER.TITLE, title) or \
            self._fill_ph(slide, PP_PLACEHOLDER.CENTER_TITLE, title)

        items = data.get("items", [])
        if not items:
            return

        step_lines = []
        for i, item in enumerate(items):
            if isinstance(item, dict):
                step_title = item.get("title", u"\u6b65\u9aa4" + str(i + 1))
                step_desc = item.get("desc", "")
            else:
                step_title = u"\u6b65\u9aa4" + str(i + 1)
                step_desc = str(item)
            step_lines.append(u"\u25b6 " + step_title)
            if step_desc:
                step_lines.append("   " + step_desc)

        body_used = self._fill_ph_multiline(slide, PP_PLACEHOLDER.BODY, step_lines)
        if not body_used:
            # Manual layout — step cards
            for i, item in enumerate(items):
                if isinstance(item, dict):
                    step_title = item.get("title", "")
                    step_desc = item.get("desc", "")
                else:
                    step_title = ""
                    step_desc = str(item)
                y = 1.4 + i * 1.0
                self._add_tb(slide, 0.7, y, 0.5, 0.5, str(i + 1), 20, True, "#FFFFFF")
                if step_title:
                    self._add_tb(slide, 1.4, y, 3.0, 0.5, step_title, 14, True)
                if step_desc:
                    self._add_tb(slide, 1.4, y + 0.4, 10.5, 0.5, step_desc, 12, False, "#555555")

    def _build_roadmap(self, data):
        idx = _find_layout_index(self.prs, "content")
        slide = self._add_slide(idx)

        title = data.get("title", "")
        self._fill_ph(slide, PP_PLACEHOLDER.TITLE, title) or \
            self._fill_ph(slide, PP_PLACEHOLDER.CENTER_TITLE, title)

        phases = data.get("items", [])
        if not phases:
            return

        for i, phase in enumerate(phases):
            if isinstance(phase, dict):
                name = phase.get("name", "")
                time = phase.get("time", "")
                milestones = phase.get("items", [])
            else:
                name = ""
                time = ""
                milestones = [str(phase)]

            x = 0.7 + i * 4.0
            header = name + (" (" + time + ")" if time else "")
            lines = [header] if header else []
            lines += [u"\u2022 " + m for m in milestones]
            self._add_multiline_tb(slide, x, 1.5, 3.5, 4.5, lines, 13)

    def _build_budget(self, data):
        idx = _find_layout_index(self.prs, "content")
        slide = self._add_slide(idx)

        title = data.get("title", "")
        self._fill_ph(slide, PP_PLACEHOLDER.TITLE, title) or \
            self._fill_ph(slide, PP_PLACEHOLDER.CENTER_TITLE, title)

        rows = data.get("items", [])
        if not rows:
            return

        lines = []
        for row in rows:
            if isinstance(row, list):
                lines.append("  |  ".join(str(c) for c in row))
            elif isinstance(row, dict):
                parts = [row.get("item", ""), row.get("amount", ""), row.get("note", "")]
                lines.append("  |  ".join(str(p) for p in parts))
            else:
                lines.append(str(row))
        self._add_multiline_tb(slide, 1.0, 1.5, 11.0, 5.0, lines, 14)

    def _build_value(self, data):
        idx = _find_layout_index(self.prs, "content")
        slide = self._add_slide(idx)

        title = data.get("title", "")
        self._fill_ph(slide, PP_PLACEHOLDER.TITLE, title) or \
            self._fill_ph(slide, PP_PLACEHOLDER.CENTER_TITLE, title)

        items = data.get("items", [])
        if not items:
            return

        count = len(items)
        if count <= 3:
            cols = count
        elif count <= 4:
            cols = 2
        else:
            cols = 3
        rows = (count + cols - 1) // cols
        area_left = 0.8
        area_top = 1.6
        area_w = 11.7
        area_h = 5.2
        gap = 0.25
        card_w = (area_w - gap * (cols - 1)) / cols
        card_h = (area_h - gap * (rows - 1)) / rows

        for i, item in enumerate(items):
            if isinstance(item, dict):
                number = item.get("number", "")
                label = item.get("label", "")
                desc = item.get("desc", "")
            else:
                number = str(item)
                label = ""
                desc = ""

            row = i // cols
            col = i % cols
            x = area_left + (card_w + gap) * col
            y = area_top + (card_h + gap) * row
            number_font = _value_font_size(number, 44 if count <= 3 else 36)

            self._add_tb(slide, x, y + 0.25, card_w, 0.8, str(number),
                         number_font, True, "#4472C4",
                         PP_ALIGN.CENTER)
            if label:
                self._add_tb(slide, x, y + 1.2, card_w, 0.4, str(label),
                             15 if count <= 3 else 13, True,
                             alignment=PP_ALIGN.CENTER)
            if desc:
                self._add_tb(slide, x, y + 1.7, card_w, max(0.7, card_h - 1.9),
                             str(desc), 11 if count <= 3 else 9, False, "#888888",
                             PP_ALIGN.CENTER)

    def _build_next_steps(self, data):
        self._build_content(data)

    def _build_guarantee(self, data):
        self._build_content(data)

    def _build_company(self, data):
        self._build_content(data)

    def _build_pricing(self, data):
        self._build_budget(data)

    def _build_delivery(self, data):
        self._build_roadmap(data)

    def _build_content(self, data):
        """Generic fallback for unknown page types."""
        idx = _find_layout_index(self.prs, "content")
        slide = self._add_slide(idx)

        title = data.get("title", "")
        subtitle = data.get("subtitle", "")
        items = data.get("items", [])

        t_ok = self._fill_ph(slide, PP_PLACEHOLDER.TITLE, title)
        if not t_ok:
            t_ok = self._fill_ph(slide, PP_PLACEHOLDER.CENTER_TITLE, title)
        if not t_ok and title:
            self._add_tb(slide, 0.7, 0.3, 11.9, 0.7, title, 28, True)

        if subtitle:
            self._add_tb(slide, 0.7, 1.0, 11.9, 0.4, subtitle, 14, False, "#888888")

        if items:
            body_used = False
            for shape in slide.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                    tf = shape.text_frame
                    tf.clear()
                    for i, item in enumerate(items):
                        if isinstance(item, dict):
                            text = item.get("name", item.get("title", item.get("desc", "")))
                        else:
                            text = str(item)
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()
                        p.text = text
                        if isinstance(item, dict) and item.get("desc") and len(text) < 50:
                            p2 = tf.add_paragraph()
                            p2.text = "    " + item["desc"]
                            p2.font.size = Pt(11)
                    body_used = True
                    break

            if not body_used:
                lines = []
                for item in items:
                    if isinstance(item, dict):
                        name = item.get("name", item.get("title", ""))
                        desc = item.get("desc", "")
                        if name:
                            lines.append(u"\u2022 " + name + (": " + desc if desc else ""))
                        elif desc:
                            lines.append(u"\u2022 " + desc)
                    elif isinstance(item, str):
                        lines.append(u"\u2022 " + item)
                self._add_multiline_tb(slide, 0.7, 1.4, 11.9, 5.0, lines, 13)

    def save(self, output_path: str):
        self.prs.save(output_path)
        return output_path
