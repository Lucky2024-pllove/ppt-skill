"""
ppt-skill CLI — 产品方案提案PPT生成器

两种生成模式:
    (A) 模板填充: python cli.py create --template brand.pptx --slides '[...]' --output out.pptx
        保留模板母版/背景/版式，只替换内容
    (B) 主题生成: python cli.py create --theme blue --slides '[...]' --output out.pptx
        使用内置主题从零构建

辅助命令:
    python cli.py list-themes
    python cli.py list-slots
"""

import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ─── 布局常量 ──────────────────────────────────────────

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_L = 0.7
MARGIN_R = 0.7
MARGIN_T = 0.5
MARGIN_B = 0.5
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
CONTENT_H = SLIDE_H - MARGIN_T - MARGIN_B

# ─── 预设主题数据 ──────────────────────────────────────

THEMES = {
    "default": {
        "label": "默认商务蓝",
        "colors": {
            "primary": "#1E2761", "secondary": "#CADCFC",
            "accent": "#F96167", "background": "#FFFFFF",
            "text": "#333333", "text_light": "#666666",
        },
        "fonts": {"title": "Calibri", "body": "Calibri", "title_ea": "微软雅黑", "body_ea": "微软雅黑"},
        "slide_size": {"width": 13.333, "height": 7.5},
    },
    "blue": {
        "label": "科技蓝",
        "colors": {
            "primary": "#0066CC", "secondary": "#E8F4FD",
            "accent": "#FF6B35", "background": "#FFFFFF",
            "text": "#333333", "text_light": "#666666",
        },
        "fonts": {"title": "Arial", "body": "Calibri", "title_ea": "微软雅黑", "body_ea": "微软雅黑"},
        "slide_size": {"width": 13.333, "height": 7.5},
    },
    "green": {
        "label": "森林绿",
        "colors": {
            "primary": "#2C5F2D", "secondary": "#97BC62",
            "accent": "#D4A017", "background": "#F5F5F5",
            "text": "#333333", "text_light": "#666666",
        },
        "fonts": {"title": "Georgia", "body": "Calibri", "title_ea": "微软雅黑", "body_ea": "微软雅黑"},
        "slide_size": {"width": 13.333, "height": 7.5},
    },
    "orange": {
        "label": "暖橙商务",
        "colors": {
            "primary": "#E8611A", "secondary": "#FDE8D8",
            "accent": "#2D3436", "background": "#FFFFFF",
            "text": "#333333", "text_light": "#666666",
        },
        "fonts": {"title": "Arial Black", "body": "Calibri", "title_ea": "微软雅黑", "body_ea": "微软雅黑"},
        "slide_size": {"width": 13.333, "height": 7.5},
    },
    "purple": {
        "label": "典雅紫",
        "colors": {
            "primary": "#6C3483", "secondary": "#E8DAEF",
            "accent": "#F1C40F", "background": "#FFFFFF",
            "text": "#333333", "text_light": "#666666",
        },
        "fonts": {"title": "Cambria", "body": "Calibri", "title_ea": "微软雅黑", "body_ea": "微软雅黑"},
        "slide_size": {"width": 13.333, "height": 7.5},
    },
    "red": {
        "label": "中国红",
        "colors": {
            "primary": "#C0392B", "secondary": "#FADBD8",
            "accent": "#F39C12", "background": "#FDFEFE",
            "text": "#333333", "text_light": "#666666",
        },
        "fonts": {"title": "Impact", "body": "Calibri", "title_ea": "微软雅黑", "body_ea": "微软雅黑"},
        "slide_size": {"width": 13.333, "height": 7.5},
    },
    "teal": {
        "label": "青碧色",
        "colors": {
            "primary": "#028090", "secondary": "#00A896",
            "accent": "#02C39A", "background": "#FDFDFD",
            "text": "#333333", "text_light": "#666666",
        },
        "fonts": {"title": "Trebuchet MS", "body": "Calibri", "title_ea": "微软雅黑", "body_ea": "微软雅黑"},
        "slide_size": {"width": 13.333, "height": 7.5},
    },
}

# ─── 页面类型插槽定义 ────────────────────────────────
# 每类页面上有哪些图片插槽（slot id），供 --images 使用

PAGE_SLOTS = {
    "cover":        ["cover_bg", "cover_logo"],
    "problem":      [],
    "solution":     ["solution_diagram"],
    "capabilities": ["capability_{i}"],  # {i} 替换为卡片序号
    "how_it_works": ["flow_diagram", "step_icon_{i}"],
    "roadmap":      ["roadmap_icon"],
    "budget":       [],
    "value":        [],
    "next_steps":   [],
    "guarantee":    ["guarantee_{i}"],
    "delivery":     [],
    "pricing":      [],
    "company":      ["company_logo", "company_photo"],
}


# ─── 颜色工具 ─────────────────────────────────────────

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def value_font_size(text: str, max_size: int = 54) -> int:
    """Choose a readable big-number font size that avoids wrapping into labels."""
    length = len(str(text or "").replace(" ", ""))
    if length <= 4:
        return max_size
    if length <= 7:
        return min(max_size, 44)
    if length <= 10:
        return min(max_size, 34)
    if length <= 14:
        return min(max_size, 28)
    return 22


# ─── 字体工具 ─────────────────────────────────────────

def _set_ea_font(run, ea_font="微软雅黑"):
    """为 run 设置东亚字体（中文显示优化）。"""
    if not ea_font:
        return
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', ea_font)


# ─── 幻灯片构建器 ──────────────────────────────────────

class SlideBuilder:
    """构建单个幻灯片，包含常用元素的辅助方法。"""

    def __init__(self, slide, theme: dict, images: dict = None):
        self.slide = slide
        self.colors = theme.get("colors", {})
        self.fonts = theme.get("fonts", {})
        self.images = images or {}
        self._add_white_background()

    def _add_white_background(self):
        bg = self.slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(self.colors.get("background", "#FFFFFF"))

    def _apply_font(self, run, font_type="body"):
        """应用字体：同时设置 Latin 与 East Asian 字体，避免中文竖排。"""
        latin_font = self.fonts.get(font_type, self.fonts.get("body", "Calibri"))
        ea_font = self.fonts.get(f"{font_type}_ea", self.fonts.get("body_ea", "微软雅黑"))
        run.font.name = latin_font
        _set_ea_font(run, ea_font)

    def add_title(self, text: str, left=None, top=None, width=None, height=None,
                  font_size=28, bold=True, color=None, alignment=PP_ALIGN.LEFT):
        left = left or Inches(MARGIN_L)
        top = top or Inches(MARGIN_T)
        width = width or Inches(CONTENT_W)
        height = height or Inches(0.7)
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = hex_to_rgb(color or self.colors.get("primary", "#333333"))
        p.alignment = alignment
        if p.runs:
            self._apply_font(p.runs[0], "title")
        return txBox

    def add_subtitle(self, text: str, font_size=14, color=None):
        txBox = self.slide.shapes.add_textbox(
            Inches(MARGIN_L), Inches(MARGIN_T + 0.8),
            Inches(CONTENT_W), Inches(0.4),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = hex_to_rgb(color or self.colors.get("text_light", "#666666"))
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            self._apply_font(p.runs[0], "body")
        return txBox

    def add_body_text(self, text: str, left=None, top=None, width=None, height=None,
                      font_size=14, color=None, bold=False, alignment=PP_ALIGN.LEFT):
        left = left or Inches(MARGIN_L)
        top = top or Inches(1.5)
        width = width or Inches(CONTENT_W)
        height = height or Inches(0.5)
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = hex_to_rgb(color or self.colors.get("text", "#333333"))
        p.font.bold = bold
        p.alignment = alignment
        if p.runs:
            self._apply_font(p.runs[0], "body")
        return txBox

    def add_bullet_list(self, items: list, left=None, top=None, width=None, height=None,
                        font_size=14, color=None):
        left = left or Inches(MARGIN_L)
        top = top or Inches(1.5)
        width = width or Inches(CONTENT_W)
        height = height or Inches(4.5)
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = hex_to_rgb(color or self.colors.get("text", "#333333"))
            p.space_after = Pt(8)
            if p.runs:
                self._apply_font(p.runs[0], "body")
            pPr = p._pPr
            if pPr is None:
                pPr = etree.SubElement(p._p, qn('a:pPr'))
            buChar = pPr.find(qn('a:buChar'))
            if buChar is None:
                buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '-')
        return txBox

    def add_rect(self, left, top, width, height, fill_color=None, line_color=None):
        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height,
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = hex_to_rgb(line_color)
        else:
            shape.line.fill.background()
        return shape

    def add_card(self, title: str, desc: str, left, top, width, height,
                 card_color=None):
        card = self.add_rect(left, top, width, height,
                             fill_color=card_color or self.colors.get("secondary", "#F0F0F0"))
        txBox = self.slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.1),
            width - Inches(0.3), Inches(0.4),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self.colors.get("primary", "#333333"))
        if p.runs:
            self._apply_font(p.runs[0], "title")
        txBox2 = self.slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.5),
            width - Inches(0.3), height - Inches(0.6),
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = hex_to_rgb(self.colors.get("text_light", "#666666"))
        if p2.runs:
            self._apply_font(p2.runs[0], "body")
        return card

    def add_image(self, slot_id: str, left, top, width=None, height=None):
        """如果 images 中有该 slot_id 对应的图片路径，则嵌入图片。"""
        img_path = self.images.get(slot_id)
        if img_path and os.path.isfile(img_path):
            w = width or Inches(4.0)
            h = height or Inches(3.0)
            try:
                pic = self.slide.shapes.add_picture(img_path, left, top, w, h)
                return pic
            except Exception:
                pass
        return None


# ─── PPT 构建器 ──────────────────────────────────────

class PptBuilder:
    """根据 slides JSON 和 images 构建 PPT。"""

    def __init__(self, theme: dict, slides: list, images: dict = None):
        self.theme = theme
        self.slides = slides
        self.images = images or {}
        self.prs = Presentation()
        sw = theme.get("slide_size", {}).get("width", 13.333)
        sh = theme.get("slide_size", {}).get("height", 7.5)
        self.prs.slide_width = Emu(int(sw * 914400))
        self.prs.slide_height = Emu(int(sh * 914400))

    def _new_sb(self):
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        return SlideBuilder(slide, self.theme, self.images)

    def build(self):
        for slide_data in self.slides:
            page_type = slide_data.get("type", "cover")
            method_name = f"_build_{page_type}"
            method = getattr(self, method_name, self._build_fallback)
            method(slide_data)

    def _build_fallback(self, data):
        """兜底：未知页面类型，仅输出标题和文本。"""
        sb = self._new_sb()
        title = data.get("title", "")
        if title:
            sb.add_title(title)
        items = data.get("items", [])
        if items:
            sb.add_bullet_list(items, top=Inches(1.5))

    def _build_cover(self, data):
        sb = self._new_sb()
        sb.add_rect(Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H),
                    fill_color=self.theme["colors"]["primary"])
        sb.add_image("cover_bg", Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
        sb.add_image("cover_logo", Inches(11.0), Inches(0.3), Inches(1.5), Inches(0.6))
        title = data.get("title", "")
        subtitle = data.get("subtitle", "")
        if title:
            sb.add_body_text(title, left=Inches(1.5), top=Inches(2.5),
                             width=Inches(10.3), height=Inches(1.2),
                             font_size=40, color="#FFFFFF", bold=True,
                             alignment=PP_ALIGN.CENTER)
        if subtitle:
            sb.add_body_text(subtitle, left=Inches(1.5), top=Inches(3.8),
                             width=Inches(10.3), height=Inches(0.6),
                             font_size=18, color="#CADCFC",
                             alignment=PP_ALIGN.CENTER)

    def _build_problem(self, data):
        sb = self._new_sb()
        title = data.get("title", "业务挑战")
        sb.add_title(title)
        subtitle = data.get("subtitle", "")
        if subtitle:
            sb.add_subtitle(subtitle)
        items = data.get("items", [])
        if items:
            sb.add_bullet_list(items, top=Inches(1.5 if not subtitle else 2.0))

    def _build_solution(self, data):
        sb = self._new_sb()
        title = data.get("title", "方案概述")
        sb.add_title(title)
        items = data.get("items", [])
        if items:
            col_w = Inches(5.0)
            left_x = Inches(MARGIN_L)
            sb.add_body_text("当前痛点", left=left_x, top=Inches(1.5),
                             width=col_w, height=Inches(0.4),
                             font_size=20, bold=True,
                             color=self.theme["colors"]["primary"])
            right_x = Inches(7.0)
            sb.add_rect(Inches(6.5), Inches(1.5), Inches(0.03), Inches(4.5),
                        fill_color=self.theme["colors"]["secondary"])
            sb.add_body_text("我们的方案", left=right_x, top=Inches(1.5),
                             width=col_w, height=Inches(0.4),
                             font_size=20, bold=True,
                             color=self.theme["colors"]["accent"])
            problems = []
            solutions = []
            for item in items:
                if isinstance(item, dict):
                    problems.append(item.get("problem", ""))
                    solutions.append(item.get("solution", ""))
                elif isinstance(item, str):
                    problems.append(item)
            if problems:
                sb.add_bullet_list(problems, left=left_x, top=Inches(2.0),
                                   width=col_w, height=Inches(3.0), font_size=13)
            if solutions:
                sb.add_bullet_list(solutions, left=right_x, top=Inches(2.0),
                                   width=col_w, height=Inches(3.0), font_size=13)
        sb.add_image("solution_diagram", Inches(7.0), Inches(2.0), Inches(4.5), Inches(3.0))
        tagline = data.get("tagline", "")
        if tagline:
            sb.add_rect(Inches(MARGIN_L), Inches(6.0), Inches(CONTENT_W), Inches(0.7),
                        fill_color=self.theme["colors"]["primary"])
            sb.add_body_text(tagline, left=Inches(1.0), top=Inches(6.1),
                             width=Inches(CONTENT_W - 0.6), height=Inches(0.5),
                             font_size=16, color="#FFFFFF", bold=True,
                             alignment=PP_ALIGN.CENTER)

    def _build_capabilities(self, data):
        sb = self._new_sb()
        title = data.get("title", "核心能力")
        sb.add_title(title)
        items = data.get("items", [])
        if not items:
            return

        # 判断内容长度：如果有任一项描述超过 30 字，改用列表模式
        use_list = any(
            len(item.get("desc", "")) > 30
            for item in items if isinstance(item, dict)
        )

        if use_list:
            # 列表模式：每项一行，左侧色条 + 标题 + 描述
            # 当项目数 > 4 时，采用双列布局以充分利用空间
            use_two_cols = len(items) > 4
            gap_in = 0.3 if use_two_cols else 0.0
            cols = 2 if use_two_cols else 1
            col_w_in = (CONTENT_W - gap_in) / cols
            items_per_col = (len(items) + 1) // 2 if use_two_cols else len(items)
            row_h_in = 0.9 if use_two_cols else 1.05
            start_y_in = 1.55
            bar_w_in = 0.06
            name_ratio = 0.32
            desc_ratio = 0.66
            inner_gap_in = 0.1
            text_h_in = row_h_in - 0.1

            for i, item in enumerate(items):
                if isinstance(item, dict):
                    name = item.get("name", "")
                    desc = item.get("desc", "")
                else:
                    name = str(item)
                    desc = ""

                col_idx = i // items_per_col if use_two_cols else 0
                row_idx = i % items_per_col if use_two_cols else i
                x_in = MARGIN_L + (col_w_in + gap_in) * col_idx
                y_in = start_y_in + row_h_in * row_idx
                x = Inches(x_in)
                y = Inches(y_in)

                # 左侧色条
                sb.add_rect(x, y, Inches(bar_w_in), Inches(text_h_in),
                            fill_color=self.theme["colors"]["primary"])

                # 标题（左列）
                name_w_in = max(0.8, col_w_in * name_ratio)
                name_left_in = x_in + bar_w_in + inner_gap_in
                sb.add_body_text(
                    name,
                    left=Inches(name_left_in),
                    top=y,
                    width=Inches(name_w_in),
                    height=Inches(text_h_in),
                    font_size=12,
                    bold=True,
                    color=self.theme["colors"]["primary"],
                )

                # 描述（右列）
                desc_left_in = name_left_in + name_w_in + inner_gap_in
                desc_w_in = max(1.0, x_in + col_w_in - desc_left_in - inner_gap_in)
                sb.add_body_text(
                    desc,
                    left=Inches(desc_left_in),
                    top=y,
                    width=Inches(desc_w_in),
                    height=Inches(text_h_in),
                    font_size=11,
                    color=self.theme["colors"]["text"],
                )
        else:
            # 网格卡片模式：短内容适用
            cols = 3
            card_w_in = (CONTENT_W - 0.4 * (cols - 1)) / cols
            card_h_in = 2.2
            start_y_in = 1.6
            gap_in = 0.4
            for i, item in enumerate(items):
                col = i % cols
                row = i // cols
                x_in = MARGIN_L + (card_w_in + gap_in) * col
                y_in = start_y_in + (card_h_in + 0.3) * row
                x = Inches(x_in)
                y = Inches(y_in)
                if isinstance(item, dict):
                    name = item.get("name", "")
                    desc = item.get("desc", "")
                else:
                    name = str(item)
                    desc = ""
                sb.add_card(name, desc, left=x, top=y,
                            width=Inches(card_w_in), height=Inches(card_h_in))
                sb.add_image(f"capability_{i+1}",
                             x + Inches(card_w_in - 1.0),
                             y + Inches(0.1),
                             Inches(0.8), Inches(0.8))

    def _build_how_it_works(self, data):
        sb = self._new_sb()
        title = data.get("title", "如何实现")
        sb.add_title(title)
        sb.add_image("flow_diagram", Inches(MARGIN_L), Inches(1.5),
                      Inches(CONTENT_W), Inches(0.3))
        items = data.get("items", [])
        if items:
            # 动态调整行高，避免项目多时溢出
            row_h_in = min(0.9, 4.8 / max(len(items), 1))
            start_y_in = 1.6
            for i, item in enumerate(items):
                if isinstance(item, dict):
                    step_title = item.get("title", f"步骤{i+1}")
                    step_desc = item.get("desc", "")
                else:
                    step_title = f"步骤{i+1}"
                    step_desc = str(item)
                y_in = start_y_in + row_h_in * i
                y = Inches(y_in)
                sb.add_rect(Inches(MARGIN_L), y, Inches(0.5), Inches(0.5),
                            fill_color=self.theme["colors"]["primary"])
                sb.add_body_text(str(i + 1), left=Inches(MARGIN_L), top=y + Inches(0.05),
                                 width=Inches(0.5), height=Inches(0.4),
                                 font_size=18, color="#FFFFFF", bold=True,
                                 alignment=PP_ALIGN.CENTER)
                sb.add_image(f"step_icon_{i+1}", Inches(MARGIN_L), y,
                             Inches(0.5), Inches(0.5))
                sb.add_body_text(step_title, left=Inches(MARGIN_L + 0.7), top=y,
                                 width=Inches(2.0), height=Inches(0.5),
                                 font_size=16, bold=True,
                                 color=self.theme["colors"]["primary"])
                sb.add_body_text(step_desc, left=Inches(MARGIN_L + 2.8), top=y,
                                 width=Inches(CONTENT_W - 3.0), height=Inches(row_h_in - 0.05),
                                 font_size=13)

    def _build_roadmap(self, data):
        sb = self._new_sb()
        title = data.get("title", "实施路线图")
        sb.add_title(title)
        sb.add_image("roadmap_icon", Inches(0.5), Inches(6.5),
                      Inches(CONTENT_W), Inches(0.3))
        phases = data.get("items", [])
        if phases:
            col_w = Inches(3.4)
            gap = Inches(0.3)
            start_x = Inches(MARGIN_L) + Inches(0.3)
            phase_colors = [
                self.theme["colors"]["primary"],
                self.theme["colors"]["accent"],
                self.theme["colors"]["primary"],
            ]
            for i, phase in enumerate(phases):
                if isinstance(phase, dict):
                    name = phase.get("name", f"Phase {i+1}")
                    time = phase.get("time", "")
                    milestones = phase.get("items", [])
                else:
                    name = f"Phase {i+1}"
                    time = ""
                    milestones = [str(phase)]
                x = start_x + (col_w + gap) * i
                c = phase_colors[i % len(phase_colors)]
                sb.add_rect(x, Inches(1.5), col_w, Inches(0.5), fill_color=c)
                header_text = f"{name}" + (f"：{time}" if time else "")
                sb.add_body_text(header_text, left=x, top=Inches(1.5),
                                 width=col_w, height=Inches(0.5),
                                 font_size=14, color="#FFFFFF", bold=True,
                                 alignment=PP_ALIGN.CENTER)
                milestone_list = [f"- {m}" for m in milestones]
                sb.add_bullet_list(milestone_list, left=x + Inches(0.1),
                                   top=Inches(2.2), width=col_w - Inches(0.2),
                                   height=Inches(2.0), font_size=13)

    def _build_budget(self, data):
        sb = self._new_sb()
        title = data.get("title", "投资预算")
        sb.add_title(title)
        rows = data.get("items", [])
        if rows:
            table_left = Inches(1.0)
            table_top = Inches(1.8)
            row_h = Inches(0.6)
            col_ws = [3.0, 3.0, 5.0]
            for i, row in enumerate(rows):
                y = table_top + row_h * i
                is_header = (i == 0)
                bg_color = self.theme["colors"]["primary"] if is_header else None
                text_color = "#FFFFFF" if is_header else self.theme["colors"]["text"]
                if isinstance(row, list):
                    cells = row
                elif isinstance(row, dict):
                    cells = [row.get("item", ""), row.get("amount", ""), row.get("note", "")]
                else:
                    cells = [str(row)]
                for j, cell_text in enumerate(cells[:3]):
                    x = table_left + sum(col_ws[:j]) * 914400
                    cw = Inches(col_ws[j])
                    if bg_color:
                        sb.add_rect(x, y, cw, row_h, fill_color=bg_color)
                    sb.add_body_text(str(cell_text), left=x + Inches(0.1),
                                     top=y + Inches(0.1),
                                     width=cw - Inches(0.2),
                                     height=row_h - Inches(0.2),
                                     font_size=14, color=text_color,
                                     bold=is_header,
                                     alignment=PP_ALIGN.CENTER if j < 2 else PP_ALIGN.LEFT)

    def _build_value(self, data):
        sb = self._new_sb()
        title = data.get("title", "价值回报")
        sb.add_title(title)
        items = data.get("items", [])
        if items:
            count = len(items)
            if count <= 3:
                cols = count
            elif count <= 4:
                cols = 2
            else:
                cols = 3
            rows = (count + cols - 1) // cols
            gap = Inches(0.25)
            area_left = Inches(MARGIN_L)
            area_top = Inches(1.65)
            area_w = Inches(CONTENT_W)
            area_h = Inches(5.25)
            card_w = int((area_w - gap * (cols - 1)) / cols)
            card_h = int((area_h - gap * (rows - 1)) / rows)

            for i, item in enumerate(items):
                if isinstance(item, dict):
                    number = item.get("number", "")
                    label = item.get("label", "")
                    desc = item.get("desc", "")
                else:
                    number = str(item)
                    label = ""
                    desc = ""

                row = i // cols
                col = i % cols
                x = area_left + (card_w + gap) * col
                y = area_top + (card_h + gap) * row

                sb.add_rect(x, y, card_w, card_h,
                            fill_color=self.theme["colors"]["secondary"])

                number_font = value_font_size(number, 54 if count <= 3 else 40)
                sb.add_body_text(str(number), left=x + Inches(0.18),
                                 top=y + Inches(0.25),
                                 width=card_w - Inches(0.36),
                                 height=Inches(0.85),
                                 font_size=number_font,
                                 color=self.theme["colors"]["primary"],
                                 bold=True, alignment=PP_ALIGN.CENTER)
                if label:
                    sb.add_body_text(str(label), left=x + Inches(0.18),
                                     top=y + Inches(1.25),
                                     width=card_w - Inches(0.36),
                                     height=Inches(0.45),
                                     font_size=16 if count <= 3 else 14,
                                     bold=True,
                                     color=self.theme["colors"]["text"],
                                     alignment=PP_ALIGN.CENTER)
                if desc:
                    sb.add_body_text(str(desc), left=x + Inches(0.22),
                                     top=y + Inches(1.78),
                                     width=card_w - Inches(0.44),
                                     height=max(Inches(0.75), card_h - Inches(2.0)),
                                     font_size=12 if count <= 3 else 10,
                                     color=self.theme["colors"]["text_light"],
                                     alignment=PP_ALIGN.CENTER)

    def _build_next_steps(self, data):
        sb = self._new_sb()
        sb.add_rect(Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H),
                    fill_color=self.theme["colors"]["secondary"])
        title = data.get("title", "下一步")
        sb.add_body_text(title, left=Inches(MARGIN_L), top=Inches(1.5),
                         width=Inches(CONTENT_W), height=Inches(0.7),
                         font_size=36, bold=True,
                         color=self.theme["colors"]["primary"],
                         alignment=PP_ALIGN.CENTER)
        items = data.get("items", [])
        if items:
            sb.add_bullet_list(items, left=Inches(3.5), top=Inches(2.5),
                               width=Inches(6.0), height=Inches(3.0),
                               font_size=16)
        contact = data.get("contact", "")
        if contact:
            sb.add_body_text(f"联系方式：{contact}", left=Inches(3.5), top=Inches(5.5),
                             width=Inches(6.0), height=Inches(0.5), font_size=14,
                             color=self.theme["colors"]["primary"], bold=True,
                             alignment=PP_ALIGN.CENTER)

    def _build_guarantee(self, data):
        sb = self._new_sb()
        title = data.get("title", "")
        if title:
            sb.add_title(title)
        items = data.get("items", [])
        if items:
            for i, item in enumerate(items):
                if isinstance(item, dict):
                    item_title = item.get("title", "")
                    item_desc = item.get("desc", "")
                else:
                    item_title = str(item)
                    item_desc = ""
                y = Inches(1.6) + Inches(1.1) * i
                sb.add_rect(Inches(MARGIN_L), y, Inches(0.08), Inches(0.6),
                            fill_color=self.theme["colors"]["primary"])
                sb.add_body_text(item_title, left=Inches(MARGIN_L + 0.3), top=y,
                                 width=Inches(2.5), height=Inches(0.6),
                                 font_size=16, bold=True,
                                 color=self.theme["colors"]["primary"])
                sb.add_body_text(item_desc, left=Inches(MARGIN_L + 3.0), top=y,
                                 width=Inches(CONTENT_W - 3.3), height=Inches(0.6),
                                 font_size=14,
                                 color=self.theme["colors"]["text"])
                sb.add_image(f"guarantee_{i+1}", Inches(MARGIN_L + 2.7), y,
                             Inches(0.3), Inches(0.3))

    def _build_delivery(self, data):
        self._build_roadmap(data)

    def _build_pricing(self, data):
        self._build_budget(data)

    def _build_company(self, data):
        sb = self._new_sb()
        title = data.get("title", "关于我们")
        sb.add_title(title)
        sb.add_image("company_logo", Inches(0.7), Inches(1.5), Inches(2.0), Inches(0.6))
        sb.add_image("company_photo", Inches(CONTENT_W - 3.0), Inches(1.5),
                      Inches(3.0), Inches(2.0))
        items = data.get("items", [])
        if items:
            sb.add_bullet_list(items, top=Inches(1.5), font_size=14)

    def save(self, output_path: str):
        self.prs.save(output_path)
        return output_path


# ─── CLI 入口 ────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="产品方案提案PPT生成器",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python cli.py create --slides 'slides.json' --output out.pptx
  python cli.py create --theme blue --slides '[...]' --images '{}' --output out.pptx
  python cli.py create --template brand.pptx --slides '[...]' --output out.pptx
  python cli.py list-themes

slides JSON 格式:
  [
    {"type":"cover","title":"产品名","subtitle":"副标题"},
    {"type":"problem","title":"业务挑战","items":["痛点1","痛点2"]},
    {"type":"solution","title":"方案概述","items":[
      {"problem":"痛点A","solution":"方案A"},
      {"problem":"痛点B","solution":"方案B"}
    ],"tagline":"一句话价值主张"},
    {"type":"capabilities","title":"核心能力","items":[
      {"name":"能力A","desc":"说明"},
      {"name":"能力B","desc":"说明"}
    ]},
    {"type":"how_it_works","title":"如何实现","items":[
      {"title":"步骤1","desc":"说明"},
      {"title":"步骤2","desc":"说明"}
    ]},
    {"type":"roadmap","title":"实施路线","items":[
      {"name":"Phase1","time":"1-2月","items":["里程碑1","里程碑2"]}
    ]},
    {"type":"budget","title":"预算","items":[
      ["项目","金额","说明"],
      ["研发","50万","内容"]
    ]},
    {"type":"value","title":"价值回报","items":[
      {"number":"80%","label":"效率","desc":"说明"}
    ]},
    {"type":"next_steps","title":"下一步","items":["动作1"],"contact":"联系方式"},
    {"type":"guarantee","title":"保障体系","items":[
      {"title":"品控保障","desc":"说明"}
    ]},
    {"type":"company","title":"关于我们","items":["公司介绍1","公司介绍2"]}
  ]
        """,
    )

    subparsers = parser.add_subparsers(dest="command", help="可用命令")

    create_parser = subparsers.add_parser("create", help="生成PPT")
    create_parser.add_argument("--template", type=str, help="模板 .pptx 文件路径")
    create_parser.add_argument("--theme", type=str, default="default",
                               choices=list(THEMES.keys()),
                               help="主题色")
    create_parser.add_argument("--slides", type=str, required=True,
                               help="幻灯片数据，JSON 字符串或 JSON 文件路径")
    create_parser.add_argument("--images", type=str, default="{}",
                               help="图片映射 JSON，如 '{\"cover_bg\":\"/path/img.png\"}'")
    create_parser.add_argument("--output", type=str, default="output.pptx",
                               help="输出文件路径")

    subparsers.add_parser("list-themes", help="查看可用主题")
    subparsers.add_parser("list-slots", help="查看所有图片插槽定义")

    args = parser.parse_args()

    if args.command == "list-themes":
        print("可用主题:")
        for key, theme in THEMES.items():
            print(f"  {key:12s} - {theme['label']}  ({theme['colors']['primary']})")
        return

    if args.command == "list-slots":
        print("图片插槽定义（每类页面上可插图的位置）:")
        for page_type, slots in PAGE_SLOTS.items():
            slot_str = ", ".join(slots) if slots else "(无)"
            print(f"  {page_type:15s}: {slot_str}")
        print()
        print("调用方Agent策略：有生图能力则生成配图传入，无生图能力则传空 {}")
        print("使用方式: --images '{\"cover_bg\":\"/path/img.png\",\"capability_1\":\"/path/img2.png\"}'")
        return

    if args.command != "create":
        parser.print_help()
        return

    # 解析 slides
    slides_data = None
    if os.path.isfile(args.slides):
        with open(args.slides, encoding="utf-8") as f:
            slides_data = json.load(f)
    else:
        slides_data = json.loads(args.slides)
    if not isinstance(slides_data, list):
        print("[ERROR] --slides 必须是一个 JSON 数组")
        sys.exit(1)

    # 解析 images
    images = {}
    if args.images:
        if os.path.isfile(args.images):
            with open(args.images, encoding="utf-8") as f:
                images = json.load(f)
        else:
            images = json.loads(args.images)

    output_path = args.output
    if not output_path.endswith(".pptx"):
        output_path += ".pptx"

    if args.template:
        # 模板填充模式：保留母版/背景/版式，只替换内容
        print(f"[OK] 基于模板填充: {args.template}")
        from template_fill import TemplatePptBuilder
        builder = TemplatePptBuilder(
            template_path=args.template,
            slides_data=slides_data,
            images=images,
        )
        builder.build()
        builder.save(output_path)
        page_count = len(builder.prs.slides)
        print(f"[OK] 生成完成: {output_path}")
        print(f"  共 {page_count} 页")
        print(f"  已保留模板母版、背景及版式布局")
    else:
        theme = THEMES.get(args.theme, THEMES["default"]).copy()
        print(f"[OK] 使用主题: {args.theme}")

        builder = PptBuilder(theme=theme, slides=slides_data, images=images)
        builder.build()
        builder.save(output_path)
        page_count = len(builder.prs.slides)
        print(f"[OK] 生成完成: {output_path}")
        print(f"  共 {page_count} 页")


if __name__ == "__main__":
    main()
